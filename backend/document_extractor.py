"""Document text extraction and chunking."""

from __future__ import annotations

import io
from concurrent.futures import ThreadPoolExecutor, as_completed

import fitz
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PDF_CHUNK_PAGES = 5
PDF_MIN_TEXT_CHARS = 20
PDF_OCR_DPI = 150
PDF_OCR_WORKERS = 4
PPTX_CHUNK_SLIDES = 5
TEXT_CHUNK_CHARS = 3000
PPTX_IMAGE_MIME = {
    "png": "image/png",
    "jpeg": "image/jpeg",
    "jpg": "image/jpeg",
    "gif": "image/gif",
    "bmp": "image/bmp",
    "tiff": "image/tiff",
    "webp": "image/webp",
}
PPTX_OCR_WORKERS = 4


def chunk_text_by_chars(text: str, chunk_size: int = TEXT_CHUNK_CHARS) -> list[str]:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    if len(normalized) <= chunk_size:
        return [normalized] if normalized.strip() else []

    chunks: list[str] = []
    start = 0
    while start < len(normalized):
        end = min(start + chunk_size, len(normalized))
        if end < len(normalized):
            split_at = normalized.rfind("\n\n", start, end)
            if split_at == -1:
                split_at = normalized.rfind("\n", start, end)
            if split_at != -1 and split_at > start:
                end = split_at
        part = normalized[start:end].strip()
        if part:
            chunks.append(part)
        start = end if end > start else start + chunk_size
    return chunks


def _page_needs_ocr(text: str) -> bool:
    return len(text.strip()) < PDF_MIN_TEXT_CHARS


def _render_pdf_page_png(page, dpi: int = PDF_OCR_DPI) -> bytes:
    zoom = dpi / 72
    matrix = fitz.Matrix(zoom, zoom)
    pixmap = page.get_pixmap(matrix=matrix, alpha=False)
    return pixmap.tobytes("png")


def _ocr_pdf_pages(
    pages_for_ocr: list[tuple[int, bytes]],
    user_id: str | None = None,
) -> dict[int, str]:
    from vision_text import extract_text_from_pdf_page_image

    results: dict[int, str] = {}
    if not pages_for_ocr:
        return results

    with ThreadPoolExecutor(max_workers=PDF_OCR_WORKERS) as executor:
        futures = {
            executor.submit(
                extract_text_from_pdf_page_image,
                image_bytes,
                page_number=page_index + 1,
                user_id=user_id,
            ): page_index
            for page_index, image_bytes in pages_for_ocr
        }
        for future in as_completed(futures):
            page_index = futures[future]
            results[page_index] = future.result()
    return results


def extract_pdf_page_texts(raw: bytes, user_id: str | None = None) -> list[str]:
    doc = fitz.open(stream=raw, filetype="pdf")
    try:
        page_texts = [page.get_text("text", sort=True) for page in doc]
        pages_for_ocr: list[tuple[int, bytes]] = []

        for page_index, text in enumerate(page_texts):
            if _page_needs_ocr(text):
                pages_for_ocr.append((page_index, _render_pdf_page_png(doc[page_index])))

        if pages_for_ocr:
            ocr_results = _ocr_pdf_pages(pages_for_ocr, user_id=user_id)
            for page_index, ocr_text in ocr_results.items():
                if ocr_text.strip():
                    page_texts[page_index] = ocr_text

        return page_texts
    finally:
        doc.close()


def _texts_from_shape(shape) -> list[str]:
    texts: list[str] = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(_texts_from_shape(child))
        return texts

    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            part = para.text.strip()
            if part:
                texts.append(part)
    else:
        try:
            part = (shape.text or "").strip()
            if part:
                texts.append(part)
        except AttributeError:
            pass

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                part = cell.text.strip()
                if part:
                    texts.append(part)

    return texts


def _images_from_shape(shape) -> list[tuple[bytes, str]]:
    images: list[tuple[bytes, str]] = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            images.extend(_images_from_shape(child))
        return images

    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return images

    try:
        image = shape.image
        blob = image.blob
        if not blob:
            return images
        ext = (image.ext or "png").lower()
        mime_type = PPTX_IMAGE_MIME.get(ext, f"image/{ext}")
        images.append((blob, mime_type))
    except Exception:
        pass

    return images


def _dedupe_lines(parts: list[str]) -> str:
    seen: set[str] = set()
    lines: list[str] = []
    for part in parts:
        for line in part.splitlines():
            cleaned = line.strip()
            if cleaned and cleaned not in seen:
                seen.add(cleaned)
                lines.append(cleaned)
    return "\n".join(lines)


def _ocr_pptx_slides(
    slides_for_ocr: list[tuple[int, list[tuple[bytes, str]]]],
    user_id: str | None = None,
) -> dict[int, str]:
    from vision_text import extract_text_from_slide_images

    results: dict[int, str] = {}
    if not slides_for_ocr:
        return results

    with ThreadPoolExecutor(max_workers=PPTX_OCR_WORKERS) as executor:
        futures = {
            executor.submit(
                extract_text_from_slide_images,
                images,
                slide_number=slide_index + 1,
                user_id=user_id,
            ): slide_index
            for slide_index, images in slides_for_ocr
        }
        for future in as_completed(futures):
            slide_index = futures[future]
            results[slide_index] = future.result()
    return results


def extract_pptx_slide_texts(raw: bytes, user_id: str | None = None) -> list[str]:
    presentation = Presentation(io.BytesIO(raw))
    slide_texts: list[str] = []
    slides_for_ocr: list[tuple[int, list[tuple[bytes, str]]]] = []

    for slide_index, slide in enumerate(presentation.slides):
        parts: list[str] = []
        images: list[tuple[bytes, str]] = []

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                note_text = notes_frame.text.strip()
                if note_text:
                    parts.append(note_text)

        for shape in slide.shapes:
            parts.extend(_texts_from_shape(shape))
            images.extend(_images_from_shape(shape))

        slide_text = _dedupe_lines(parts)
        if slide_text:
            slide_texts.append(slide_text)
        elif images:
            slide_texts.append("")
            slides_for_ocr.append((slide_index, images))
        else:
            slide_texts.append("")

    if slides_for_ocr:
        ocr_results = _ocr_pptx_slides(slides_for_ocr, user_id=user_id)
        for slide_index, ocr_text in ocr_results.items():
            slide_texts[slide_index] = ocr_text

    return slide_texts


def build_slide_chunks(slide_texts: list[str], chunk_size: int = PPTX_CHUNK_SLIDES) -> list[dict]:
    total_slides = len(slide_texts)
    return [
        {
            "text": "\n".join(slide_texts[i : i + chunk_size]),
            "start_page": i + 1,
            "end_page": min(i + chunk_size, total_slides),
        }
        for i in range(0, total_slides, chunk_size)
        if any(slide_texts[i : i + chunk_size])
    ]


def build_pdf_chunks(pdf_page_texts: list[str]) -> list[dict]:
    total_pages = len(pdf_page_texts)
    return [
        {
            "text": "\n".join(pdf_page_texts[i : i + PDF_CHUNK_PAGES]),
            "start_page": i + 1,
            "end_page": min(i + PDF_CHUNK_PAGES, total_pages),
        }
        for i in range(0, total_pages, PDF_CHUNK_PAGES)
        if any(pdf_page_texts[i : i + PDF_CHUNK_PAGES])
    ]


def extract_text_from_upload(
    filename: str,
    raw: bytes,
    user_id: str | None = None,
) -> tuple[str, list[dict]]:
    file_type = filename.rsplit(".", 1)[-1].lower()

    if file_type == "pdf":
        pdf_page_texts = extract_pdf_page_texts(raw, user_id=user_id)
        text = "\n".join(pdf_page_texts)
        source_chunks = build_pdf_chunks(pdf_page_texts)
        if not source_chunks:
            raise ValueError(
                "PDF에서 텍스트를 읽지 못했습니다. "
                "스캔 문서는 OCR을 시도했으나 내용을 찾지 못했습니다."
            )
    elif file_type == "docx":
        document = Document(io.BytesIO(raw))
        text = "\n".join(para.text for para in document.paragraphs if para.text.strip())
        source_chunks = [{"text": chunk} for chunk in chunk_text_by_chars(text)]
    elif file_type == "pptx":
        slide_texts = extract_pptx_slide_texts(raw, user_id=user_id)
        text = "\n".join(slide_texts)
        source_chunks = build_slide_chunks(slide_texts)
        if not source_chunks:
            raise ValueError(
                "슬라이드에서 텍스트를 읽지 못했습니다. "
                "이미지로만 된 슬라이드는 OCR을 시도했으나 내용을 찾지 못했습니다."
            )
    else:
        raise ValueError("PDF, DOCX, PPTX 파일만 지원합니다.")

    return text, source_chunks
